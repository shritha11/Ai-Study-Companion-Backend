from pptx import Presentation
def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    texts.append(shape.text)
    return "\n".join(texts)